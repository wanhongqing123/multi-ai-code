#!/usr/bin/env python3
"""Generate docs/Multi-AI-Code-介绍.pptx — editorial product-tour deck."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

try:
    from PIL import Image
except ImportError:
    Image = None


# ── Palette: editorial warm cream + berry accent ──
BG       = RGBColor(0xFA, 0xF6, 0xEE)   # soft cream base (content slides)
BG_DARK  = RGBColor(0x4A, 0x1E, 0x2E)   # deep berry (cover + closing)
CARD     = RGBColor(0xFF, 0xFF, 0xFF)
TINT     = RGBColor(0xF5, 0xED, 0xE0)   # richer cream for highlight fills
INK      = RGBColor(0x1C, 0x19, 0x17)   # stone 900
MUTED    = RGBColor(0x78, 0x71, 0x6C)   # stone 500
HAIRLINE = RGBColor(0xE7, 0xE5, 0xE4)   # stone 200
BERRY    = RGBColor(0x6D, 0x2E, 0x46)   # primary accent
ROSE     = RGBColor(0xA2, 0x67, 0x69)   # secondary accent
OCHRE    = RGBColor(0xC2, 0x41, 0x0C)   # sharp accent, used sparingly
CREAM    = RGBColor(0xF5, 0xED, 0xE0)   # text color on dark backgrounds
CREAM_MUTED = RGBColor(0xD4, 0xC4, 0xAB) # muted text on dark
NUM_FG   = RGBColor(0xFF, 0xFF, 0xFF)

# Fonts — Songti SC covers Chinese + English glyphs on Mac
SERIF = "Songti SC"
SANS  = "PingFang SC"
MONO  = "Menlo"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(1.0)
MARGIN_R = Inches(1.0)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R

SCREENSHOT_DIR = Path(__file__).resolve().parent.parent / "docs" / "screenshots"


# ─────────────────────── primitives ───────────────────────

def _send_to_back(shape):
    spTree = shape._element.getparent()
    spTree.remove(shape._element)
    spTree.insert(2, shape._element)


def fill_rect(slide, left, top, width, height, color, *, no_line=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if no_line:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def set_bg(slide, color=BG):
    bg = fill_rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)
    _send_to_back(bg)
    return bg


def add_text(slide, left, top, width, height, text, *,
             font=SANS, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top  = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                size=15, color=INK, line_spacing=1.5, dash_color=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top  = Emu(0); tf.margin_bottom = Emu(0)

    dc = dash_color if dash_color else BERRY
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(7)
        p.line_spacing = line_spacing
        dash = p.add_run()
        dash.text = "—   "
        dash.font.name = SANS
        dash.font.size = Pt(size)
        dash.font.color.rgb = dc
        dash.font.bold = True
        body = p.add_run()
        body.text = item
        body.font.name = SANS
        body.font.size = Pt(size)
        body.font.color.rgb = color
    return tb


def rrect(slide, left, top, width, height, *, fill=CARD,
          border_color=None, radius=0.04):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                   width, height)
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


# ─────────────────────── chrome (every slide) ───────────────────────

def chrome(slide, page_idx, total, *, dark=False):
    """Minimal editorial chrome: brand top-left, page top-right.
    No color bars / rules anywhere (skill forbids accent lines under titles)."""
    brand_color = CREAM if dark else BERRY
    page_color = CREAM_MUTED if dark else MUTED
    add_text(slide, MARGIN_L, Inches(0.4), Inches(6), Inches(0.3),
             "MULTI-AI  CODE",
             font=SANS, size=9, bold=True, color=brand_color)
    add_text(slide, SLIDE_W - MARGIN_R - Inches(2), Inches(0.4),
             Inches(2), Inches(0.3),
             f"{page_idx:02d}  ／  {total:02d}",
             font=MONO, size=9, color=page_color, align=PP_ALIGN.RIGHT)


def add_chapter_header(slide, chapter_num, kicker, title, *,
                       title_size=34, width=Inches(10),
                       chapter_color=BERRY):
    """Giant serif chapter numeral + small kicker + title.
    No horizontal rule — the numeral itself is the motif."""
    # huge numeral top-left
    add_text(slide, MARGIN_L, Inches(0.9), Inches(2.4), Inches(2.2),
             f"{chapter_num:02d}",
             font=SERIF, size=110, bold=True, color=chapter_color,
             line_spacing=1.0)
    # kicker + title, offset to the right of the numeral
    text_left = MARGIN_L + Inches(2.6)
    add_text(slide, text_left, Inches(1.25), width, Inches(0.35),
             kicker,
             font=SANS, size=11, bold=True, color=OCHRE)
    add_text(slide, text_left, Inches(1.7), width, Inches(1.6),
             title,
             font=SERIF, size=title_size, bold=True, color=INK,
             line_spacing=1.15)


# ─────────────────────── components ───────────────────────

def add_screenshot(slide, left, top, width, height, filename, *, caption=None):
    """Fit screenshot to (width, height) preserving aspect ratio.
    No outer frame padding — the image itself sits on a hairline-bordered card
    sized to its final rendered dimensions (no dead white space top/bottom)."""
    path = SCREENSHOT_DIR / filename
    if path.exists():
        if Image is not None:
            with Image.open(path) as img:
                iw, ih = img.size
            aspect = iw / ih if ih else 1.4
        else:
            aspect = 1.4
        # Fit to width first; if it overshoots allowed height, fit to height
        pic_w = width
        pic_h = int(pic_w / aspect)
        if pic_h > height:
            pic_h = height
            pic_w = int(pic_h * aspect)
        pic_left = left + (width - pic_w) // 2  # horizontal center in allowed box
        pic_top = top                            # top-align (match text block top)
        # hairline-bordered backing frame sized to the image
        rrect(slide, pic_left, pic_top, pic_w, pic_h,
              fill=CARD, border_color=HAIRLINE, radius=0.02)
        slide.shapes.add_picture(str(path), pic_left, pic_top, pic_w, pic_h)
        # adjust caption position if provided — sits directly under the image
        if caption:
            add_text(slide, pic_left, pic_top + pic_h + Inches(0.15),
                     pic_w, Inches(0.3),
                     caption,
                     font=SANS, size=10, color=MUTED,
                     align=PP_ALIGN.CENTER)
        return  # early return; skip the below-frame caption block
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       left, top, width, height)
        shape.adjustments[0] = 0.03
        shape.fill.solid()
        shape.fill.fore_color.rgb = TINT
        shape.line.color.rgb = ROSE
        shape.line.width = Pt(0.75)
        shape.line.dash_style = 7
        shape.shadow.inherit = False
        tf = shape.text_frame
        tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"截图占位\ndocs/screenshots/{filename}"
        r.font.name = MONO; r.font.size = Pt(12); r.font.color.rgb = MUTED
    if caption:
        add_text(slide, left, top + height + Inches(0.1),
                 width, Inches(0.3),
                 caption,
                 font=SANS, size=10, color=MUTED,
                 align=PP_ALIGN.CENTER)


def add_pill(slide, left, top, text, *, fill=CARD, color=BERRY):
    char_count = sum(2 if ord(c) > 255 else 1 for c in text)
    width = Inches(0.4 + char_count * 0.09)
    height = Inches(0.4)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = color
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = SANS; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = color
    return left + width


def add_icon_circle(slide, cx, cy, diameter, numeral, *,
                    fill=BERRY, color=NUM_FG):
    """Filled circle with a big serif numeral inside — used as an icon badge."""
    left = cx - diameter / 2
    top = cy - diameter / 2
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = numeral
    r.font.name = SERIF
    r.font.size = Pt(int(diameter / Inches(1) * 32))
    r.font.bold = True
    r.font.color.rgb = color


def add_pain_card(slide, left, top, width, height, *, icon, title, body):
    rrect(slide, left, top, width, height, fill=CARD, border_color=HAIRLINE)
    # icon circle on upper-left
    diameter = Inches(0.5)
    icon_cx = left + Inches(0.35) + diameter / 2
    icon_cy = top + Inches(0.35) + diameter / 2
    add_icon_circle(slide, icon_cx, icon_cy,
                    diameter, icon,
                    fill=TINT, color=BERRY)
    # title to the right of the icon, same vertical band
    add_text(slide, left + Inches(1.05), top + Inches(0.3),
             width - Inches(1.25), Inches(0.55),
             title,
             font=SERIF, size=17, bold=True, color=INK)
    # body — single line, generous width, sits in lower half
    add_text(slide, left + Inches(0.42), top + Inches(1.0),
             width - Inches(0.85), height - Inches(1.1),
             body,
             font=SANS, size=13, color=MUTED, line_spacing=1.45)


def add_scenario_card(slide, left, top, width, height, *,
                      kicker, title, body):
    rrect(slide, left, top, width, height, fill=CARD, border_color=HAIRLINE)
    add_text(slide, left + Inches(0.35), top + Inches(0.3),
             width - Inches(0.7), Inches(0.3),
             kicker,
             font=SANS, size=10, bold=True, color=OCHRE)
    add_text(slide, left + Inches(0.35), top + Inches(0.65),
             width - Inches(0.7), Inches(0.6),
             title,
             font=SERIF, size=22, bold=True, color=BERRY)
    add_text(slide, left + Inches(0.35), top + Inches(1.55),
             width - Inches(0.7), height - Inches(1.7),
             body,
             font=SANS, size=13, color=INK, line_spacing=1.55)


def build_tour_slide(prs, blank, page_idx, total, *,
                     tour_n, tour_total, kicker, title, bullets,
                     screenshot_name):
    """Two-column editorial tour: kicker + title + bullets on the left,
    screenshot on the right. Both columns top-align at the same Y."""
    s = prs.slides.add_slide(blank)
    set_bg(s)
    chrome(s, page_idx, total)

    text_w = Inches(5.3)
    gap = Inches(0.4)
    right_x = MARGIN_L + text_w + gap
    right_w = CONTENT_W - text_w - gap

    # shared top baseline for both columns
    top_y = Inches(1.4)

    # kicker — "产品导览   01 / 04"
    add_text(s, MARGIN_L, top_y, text_w, Inches(0.4),
             f"{kicker}   ／   {tour_n:02d}  ／  {tour_total:02d}",
             font=SANS, size=11, bold=True, color=OCHRE)

    # title — large serif, directly under kicker
    add_text(s, MARGIN_L, top_y + Inches(0.5), text_w, Inches(1.9),
             title,
             font=SERIF, size=26, bold=True, color=INK, line_spacing=1.25)

    # bullets — tight below title (no big gap)
    add_bullets(s, MARGIN_L, top_y + Inches(2.4), text_w, Inches(3.5),
                bullets, size=14, line_spacing=1.6)

    # screenshot — same top as kicker, filling right column naturally
    add_screenshot(s, right_x, top_y,
                   right_w, Inches(5.5), screenshot_name)
    return s


# ─────────────────────── slides ───────────────────────

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    TOTAL = 10

    def capability_card(slide, left, top, width, height, *,
                        tag, title, body):
        rrect(slide, left, top, width, height, fill=CARD, border_color=HAIRLINE)
        add_text(slide, left + Inches(0.3), top + Inches(0.24),
                 width - Inches(0.6), Inches(0.3),
                 tag,
                 font=SANS, size=10, bold=True, color=OCHRE)
        add_text(slide, left + Inches(0.3), top + Inches(0.62),
                 width - Inches(0.6), Inches(0.55),
                 title,
                 font=SERIF, size=20, bold=True, color=BERRY)
        add_text(slide, left + Inches(0.3), top + Inches(1.35),
                 width - Inches(0.6), height - Inches(1.55),
                 body,
                 font=SANS, size=13, color=INK, line_spacing=1.45)

    def simple_value_card(slide, left, top, width, height, *,
                          title, body):
        rrect(slide, left, top, width, height, fill=CARD, border_color=HAIRLINE)
        add_text(slide, left + Inches(0.3), top + Inches(0.28),
                 width - Inches(0.6), Inches(0.55),
                 title,
                 font=SERIF, size=18, bold=True, color=INK)
        add_text(slide, left + Inches(0.3), top + Inches(0.95),
                 width - Inches(0.6), height - Inches(1.15),
                 body,
                 font=SANS, size=13, color=MUTED, line_spacing=1.45)

    # Slide 1 · Cover
    s = prs.slides.add_slide(blank)
    set_bg(s, BG_DARK)
    add_text(s, MARGIN_L, Inches(0.75), Inches(5.5), Inches(0.35),
             'A LOCAL AI CODING WORKBENCH',
             font=SANS, size=11, bold=True, color=OCHRE)
    add_text(s, MARGIN_L, Inches(1.2), Inches(6.3), Inches(1.2),
             'Multi-AI',
             font=SERIF, size=62, bold=True, color=CREAM, line_spacing=0.95)
    add_text(s, MARGIN_L, Inches(2.0), Inches(6.3), Inches(1.2),
             'Code',
             font=SERIF, size=62, bold=True, color=OCHRE, line_spacing=0.95)
    add_text(s, MARGIN_L, Inches(3.15), Inches(5.4), Inches(1.0),
             '面向本地仓库的 AI CLI 协作工作台。\n把方案阅读、主会话推进、Diff 审查与反馈回灌放进一个工具里。',
             font=SERIF, size=18, color=CREAM_MUTED, line_spacing=1.45)
    x = MARGIN_L
    for tag in ['方案驱动', 'Diff 审查', 'Claude Code / Codex', '本地运行']:
        x = add_pill(s, x, Inches(5.15), tag, fill=BG_DARK, color=CREAM) + Inches(0.12)
    add_screenshot(s, Inches(7.0), Inches(1.0), Inches(5.3), Inches(5.9), '主界面.png')
    add_text(s, SLIDE_W - Inches(4.3), Inches(0.7), Inches(3.8), Inches(3.4),
             '01',
             font=SERIF, size=145, bold=True,
             color=RGBColor(0x5E, 0x2A, 0x3A),
             align=PP_ALIGN.RIGHT, line_spacing=1.0)
    chrome(s, 1, TOTAL, dark=True)

    # Slide 2 · Why now / pain points
    s = prs.slides.add_slide(blank)
    set_bg(s)
    chrome(s, 2, TOTAL)
    add_chapter_header(s, 2, '为什么需要它', '问题不在 AI 不够强，而在协作流程不够顺')
    card_w = Inches(5.4)
    card_h = Inches(1.72)
    row1_y = Inches(3.35)
    row2_y = row1_y + card_h + Inches(0.28)
    col1_x = MARGIN_L
    col2_x = MARGIN_L + card_w + Inches(0.3)
    pains = [
        (col1_x, row1_y, '①', '方案和执行脱节',
         '讨论过的方向留在聊天里，\nAI 一旦开始写代码，就容易偏离最初意图。'),
        (col2_x, row1_y, '②', '长会话难以收束',
         '上下文越长越容易漂移，\n到后面很难判断 AI 现在到底在按什么目标推进。'),
        (col1_x, row2_y, '③', 'Diff 反馈转述成本高',
         '用户能在代码里看出问题，\n但还要手动复制片段、描述位置，再重新讲给 AI。'),
        (col2_x, row2_y, '④', 'CLI 工具各自为战',
         'Claude Code 和 Codex 各有优势，\n但缺少统一工作台把它们接进同一条流程。'),
    ]
    for (x, y, ic, title, body) in pains:
        add_pain_card(s, x, y, card_w, card_h, icon=ic, title=title, body=body)
    add_text(s, MARGIN_L, Inches(6.83), CONTENT_W, Inches(0.36),
             'Multi-AI Code 解决的不是单点能力，而是让方案、执行、审查和反馈形成闭环。',
             font=SERIF, size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # Slide 3 · Product positioning
    s = prs.slides.add_slide(blank)
    set_bg(s)
    chrome(s, 3, TOTAL)
    add_chapter_header(s, 3, '产品定位', '它不是普通聊天壳，而是一个面向真实仓库的 AI 工作台')
    add_text(s, MARGIN_L, Inches(3.25), Inches(4.6), Inches(1.45),
             '一个项目 = 一个本地仓库 + 一份当前方案 + 一个主会话。\nAI 先理解方案，再推进实现；用户通过 Diff 审查把反馈继续回灌给同一个会话。',
             font=SANS, size=17, color=INK, line_spacing=1.5)
    labels = [
        ('01', '方案', '任务目标显式落在 Markdown 文件里'),
        ('02', '主会话', 'AI CLI 在真实 PTY 终端里推进任务'),
        ('03', 'Diff 审查', '用户针对具体改动逐行提出反馈'),
        ('04', '回灌闭环', '反馈直接回到当前会话，继续修改'),
    ]
    start_x = MARGIN_L
    gap = Inches(0.22)
    box_w = (CONTENT_W - gap * 3) / 4
    for idx, (num, title, body) in enumerate(labels):
        x = start_x + idx * (box_w + gap)
        rrect(s, x, Inches(5.25), box_w, Inches(1.25), fill=TINT, border_color=None)
        add_text(s, x + Inches(0.22), Inches(5.42), Inches(0.6), Inches(0.32),
                 num, font=MONO, size=10, bold=True, color=BERRY)
        add_text(s, x + Inches(0.22), Inches(5.74), box_w - Inches(0.44), Inches(0.32),
                 title, font=SERIF, size=17, bold=True, color=INK)
        add_text(s, x + Inches(0.22), Inches(6.03), box_w - Inches(0.44), Inches(0.42),
                 body, font=SANS, size=10.5, color=MUTED, line_spacing=1.25)

    # Slide 4 · Capability overview
    s = prs.slides.add_slide(blank)
    set_bg(s)
    chrome(s, 4, TOTAL)
    add_chapter_header(s, 4, '核心能力总览', '当前版本已经形成一条可使用、可演示的主流程')
    card_w = Inches(5.4)
    card_h = Inches(2.2)
    top_y = Inches(3.3)
    capability_card(s, MARGIN_L, top_y, card_w, card_h,
                    tag='01', title='主会话终端',
                    body='以真实 PTY 驱动 Claude Code 或 Codex。\n状态清晰，可启动、停止、重启，适合持续推进单个任务。')
    capability_card(s, MARGIN_L + card_w + Inches(0.3), top_y, card_w, card_h,
                    tag='02', title='方案导入与预览',
                    body='支持仓库内方案与外部方案文件。\n先看方案、先做标注、先确认方向，再决定是否继续修改代码。')
    capability_card(s, MARGIN_L, top_y + card_h + Inches(0.28), card_w, card_h,
                    tag='03', title='Diff 审查与逐行反馈',
                    body='双栏查看当前改动，支持多行选区、逐条说明与整体意见。\n反馈可直接回灌给当前会话。')
    capability_card(s, MARGIN_L + card_w + Inches(0.3), top_y + card_h + Inches(0.28), card_w, card_h,
                    tag='04', title='双 AI CLI 配置',
                    body='支持 Claude Code 与 Codex。\n团队可以按已有习惯选择 CLI，并保留 binary、参数与环境变量级别的灵活性。')

    # Slide 5 · Typical workflow
    s = prs.slides.add_slide(blank)
    set_bg(s)
    chrome(s, 5, TOTAL)
    add_chapter_header(s, 5, '典型使用流程', '从进入仓库到完成一轮迭代，路径足够清晰')
    steps = [
        ('1', '选择仓库', 'AI CLI 在真实项目目录下工作'),
        ('2', '选择方案', '使用仓库内方案，或导入外部 Markdown'),
        ('3', '启动会话', 'AI 先阅读方案，总结目标并等待确认'),
        ('4', '推进实现', '在主终端里持续修改、运行、解释'),
        ('5', '打开 Diff 审查', '查看当前代码改动，选中关键位置'),
        ('6', '批注回灌', '把逐行反馈发回同一个会话继续修改'),
    ]
    step_w = (CONTENT_W - Inches(1.0)) / 6
    base_y = Inches(4.1)
    for i, (num, title, body) in enumerate(steps):
        x = MARGIN_L + i * (step_w + Inches(0.2))
        add_icon_circle(s, x + step_w / 2, base_y, Inches(0.78), num,
                        fill=BERRY if i in (2, 4) else TINT,
                        color=NUM_FG if i in (2, 4) else BERRY)
        if i in (0, 1, 3, 5):
            slide_border_circle(s, x + step_w / 2, base_y, Inches(0.78))
        add_text(s, x + Inches(0.02), base_y + Inches(0.55), step_w - Inches(0.04), Inches(0.42),
                 title, font=SERIF, size=15, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.02), base_y + Inches(1.0), step_w - Inches(0.04), Inches(0.9),
                 body, font=SANS, size=11, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.35)
        if i < len(steps) - 1:
            fill_rect(s, x + step_w - Inches(0.02), base_y - Emu(3500),
                      Inches(0.16), Emu(7000), HAIRLINE)
    add_text(s, MARGIN_L, Inches(6.78), CONTENT_W, Inches(0.34),
             '关键点：AI 不是一次性跑完，而是在用户的确认和批注中逐轮推进。',
             font=SERIF, size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # Slide 6 · Plan-driven workflow
    s = prs.slides.add_slide(blank)
    set_bg(s)
    chrome(s, 6, TOTAL)
    add_text(s, MARGIN_L, Inches(1.0), Inches(5.4), Inches(0.3),
             '方案驱动能力',
             font=SANS, size=11, bold=True, color=OCHRE)
    add_text(s, MARGIN_L, Inches(1.42), Inches(5.6), Inches(0.9),
             '方案文件是一等公民，\n而不是沉在聊天记录里。',
             font=SERIF, size=30, bold=True, color=INK, line_spacing=1.2)
    add_bullets(s, MARGIN_L, Inches(3.0), Inches(4.8), Inches(2.5), [
        '支持仓库内方案和外部 Markdown 方案',
        '外部方案只把绝对路径交给 AI CLI，自行读取，避免长文刷屏',
        '在用户确认之前，AI 先总结方案，不直接改代码',
        '方案预览本身就支持标注，适合先对齐目标再进入实施',
    ], size=14, line_spacing=1.55)
    add_screenshot(s, Inches(6.8), Inches(1.1), Inches(5.1), Inches(5.7), '方案 Review.png')

    # Slide 7 · Diff review
    s = prs.slides.add_slide(blank)
    set_bg(s)
    chrome(s, 7, TOTAL)
    add_text(s, MARGIN_L, Inches(1.0), Inches(5.2), Inches(0.3),
             'Diff 审查与批注回灌',
             font=SANS, size=11, bold=True, color=OCHRE)
    add_text(s, MARGIN_L, Inches(1.42), Inches(5.8), Inches(0.95),
             '看代码的时候就把意见写下来，\n而不是事后再转述给 AI。',
             font=SERIF, size=29, bold=True, color=INK, line_spacing=1.2)
    add_bullets(s, MARGIN_L, Inches(3.05), Inches(4.9), Inches(2.65), [
        '查看工作区改动、最近提交或指定 commit',
        '支持单行、多行代码选区标注',
        '右侧逐条填写说明，也可补整体意见',
        '一键把批注发回当前 AI CLI，继续修改',
    ], size=14, line_spacing=1.55)
    add_screenshot(s, Inches(6.55), Inches(1.0), Inches(5.35), Inches(5.8), '代码 Review.png')

    # Slide 8 · AI CLI settings
    s = prs.slides.add_slide(blank)
    set_bg(s)
    chrome(s, 8, TOTAL)
    add_chapter_header(s, 8, 'AI CLI 配置能力', '不改变团队的 CLI 习惯，而是把它们接进统一工作流')
    capability_card(s, MARGIN_L, Inches(3.3), Inches(3.55), Inches(2.25),
                    tag='CLI', title='Claude Code',
                    body='适合强调交互式协作、阅读与修改节奏的团队。')
    capability_card(s, MARGIN_L + Inches(3.85), Inches(3.3), Inches(3.55), Inches(2.25),
                    tag='CLI', title='Codex',
                    body='适合希望直接用 Codex CLI 推进主任务，统一放在同一个桌面工具里操作。')
    add_bullets(s, MARGIN_L + Inches(7.95), Inches(3.45), Inches(3.2), Inches(2.6), [
        '支持切换 `Claude Code` / `Codex`',
        '支持 binary override',
        '支持附加参数与环境变量',
        '更适配不同团队已有的 AI CLI 习惯',
    ], size=14, line_spacing=1.55)
    add_text(s, MARGIN_L, Inches(6.58), CONTENT_W, Inches(0.42),
             '价值不在于绑定某一个模型，而在于把不同 CLI 工具收进同一个稳定工作流。',
             font=SERIF, size=15, color=MUTED, align=PP_ALIGN.CENTER)

    # Slide 9 · Product value
    s = prs.slides.add_slide(blank)
    set_bg(s)
    chrome(s, 9, TOTAL)
    add_chapter_header(s, 9, '当前产品价值', '从老板视角看，这个工具已经有明确的投入方向')
    card_w = Inches(5.4)
    card_h = Inches(1.75)
    simple_value_card(s, MARGIN_L, Inches(3.35), card_w, card_h,
                      title='让方案到实现更可控',
                      body='方案先显式落盘，AI 先阅读、先总结、先确认，再继续改代码。')
    simple_value_card(s, MARGIN_L + card_w + Inches(0.3), Inches(3.35), card_w, card_h,
                      title='让代码审查形成闭环',
                      body='用户在 Diff 审查里看到问题，就地批注，再直接发回 AI 继续改。')
    simple_value_card(s, MARGIN_L, Inches(5.38), card_w, card_h,
                      title='降低团队使用门槛',
                      body='把真实 CLI、方案、终端和审查集中到一个桌面工具里，减少切换成本。')
    simple_value_card(s, MARGIN_L + card_w + Inches(0.3), Inches(5.38), card_w, card_h,
                      title='为后续沉淀打基础',
                      body='当前已经具备演示能力，后续可继续拓展历史、记忆和协作能力。')

    # Slide 10 · Next opportunities / close
    s = prs.slides.add_slide(blank)
    set_bg(s, BG_DARK)
    chrome(s, 10, TOTAL, dark=True)
    add_text(s, MARGIN_L, Inches(1.0), CONTENT_W, Inches(0.34),
             '下一步机会',
             font=SANS, size=11, bold=True, color=OCHRE)
    add_text(s, MARGIN_L, Inches(1.45), CONTENT_W, Inches(1.0),
             '当前方向已经清晰，\n接下来重点是把体验继续打磨到更稳定、更顺手。',
             font=SERIF, size=34, bold=True, color=CREAM, line_spacing=1.15)
    next_items = [
        ('01', '更稳定的长会话体验'),
        ('02', '更强的 markdown 与输出可读性'),
        ('03', '更完善的历史、上下文与项目级沉淀'),
        ('04', '更成熟的团队演示与协作能力'),
    ]
    row_y = Inches(4.1)
    item_w = (CONTENT_W - Inches(0.9)) / 2
    item_h = Inches(0.92)
    for idx, (num, title) in enumerate(next_items):
        col = idx % 2
        row = idx // 2
        x = MARGIN_L + col * (item_w + Inches(0.3))
        y = row_y + row * (item_h + Inches(0.28))
        rrect(s, x, y, item_w, item_h, fill=RGBColor(0x5A, 0x25, 0x37), border_color=None)
        add_text(s, x + Inches(0.22), y + Inches(0.23), Inches(0.5), Inches(0.22),
                 num, font=MONO, size=10, bold=True, color=CREAM_MUTED)
        add_text(s, x + Inches(0.75), y + Inches(0.18), item_w - Inches(0.95), Inches(0.32),
                 title, font=SERIF, size=18, bold=True, color=CREAM)
    add_text(s, MARGIN_L, Inches(6.78), CONTENT_W, Inches(0.34),
             'Multi-AI Code 已具备明确主线与可演示能力，适合继续打磨为更成熟的 AI 编码工作台。',
             font=SERIF, size=16, color=CREAM_MUTED)

    out = Path(__file__).resolve().parent.parent / 'docs' / 'Multi-AI-Code-老板展示版.pptx'
    prs.save(out)
    print(f'Wrote {out}')

def slide_border_circle(slide, cx, cy, diameter):
    """Overlay a hollow circle on top of a filled one, used to give
    non-focus step circles a thin berry border."""
    left = cx - diameter / 2
    top = cy - diameter / 2
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shape.fill.background()
    shape.line.color.rgb = BERRY
    shape.line.width = Pt(1.2)
    shape.shadow.inherit = False
    return shape


if __name__ == "__main__":
    build()
